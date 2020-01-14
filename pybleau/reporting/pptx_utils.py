""" Basic tools to simplify building pptx files from data.
"""
import logging

try:
    from pptx import Presentation
    from pptx.util import Inches
except ImportError:
    msg = "Missing the python-pptx package this utility is relying on."
    raise ImportError(msg)

logger = logging.getLogger(__name__)


def title_slide(presentation=None, title_text="", sub_title_text=""):
    """ Add title slide to specified presentation.

    Parameters
    ----------
    presentation : Presentation, optional
        Presentation to add the slide to. If left as None, a new presentation
        is created and returned.

    title_text : str, optional
        Title of the slide, if any.

    sub_title_text : str, optional
        Sub-title of the slide, if any.
    """
    if presentation is None:
        presentation = Presentation()

    title_slide_layout = presentation.slide_layouts[0]
    slide = presentation.slides.add_slide(title_slide_layout)

    if title_text:
        title = slide.shapes.title
        title.text = title_text

    if sub_title_text:
        subtitle = slide.placeholders[1]
        subtitle.text = sub_title_text

    return presentation


def df_to_slide(presentation, data, slide_title="", left=None, top=None,
                width=None, height=None, include_column_names=True,
                include_index=False):
    """ Add slide to specified presentation with table containing the data.

    Parameters
    ----------
    presentation : Presentation
        Presentation to add the slide to.

    data : pd.DataFrame
        DataFrame containing the data to display in the new slide's table.

    slide_title : str, optional
        Title of the slide, if any.

    left : Inches, float, int, optional
        Distance (in inches) from the left side of the new slide and the left
        side of the image.

    top : Inches, float, int, optional
        Distance (in inches) from the top of the new slide and the top of the
        image.

    width : Inches, float, int, optional
        Width (in inches) of the table element.

    height : Inches, float, int, optional
        Height (in inches) of the table element.

    include_column_names : bool, optional
        Whether to include a row for the column names.

    include_index : bool, optional
        Whether to include a column for the index values.
    """
    from pptx.presentation import Presentation

    if not isinstance(presentation, Presentation):
        msg = "The presentation should be a python-pptx Presentation."
        logger.exception(msg)
        raise ValueError(msg)

    title_only_slide_layout = presentation.slide_layouts[5]
    slide = presentation.slides.add_slide(title_only_slide_layout)
    shapes = slide.shapes

    if slide_title:
        title_shape = shapes.title
        title_shape.text = slide_title

    if left is None:
        left = Inches(1.2)
    elif isinstance(left, (float, int)):
        left = Inches(left)

    if top is None:
        top = Inches(2.0)
    elif isinstance(top, (float, int)):
        top = Inches(top)

    if width is None:
        width = Inches(6.0)
    elif isinstance(width, (float, int)):
        width = Inches(width)

    if height is None:
        height = Inches(0.8)
    elif isinstance(height, (float, int)):
        height = Inches(height)

    num_rows = data.shape[0]
    if include_column_names:
        num_rows += 1

    num_cols = data.shape[1]
    if include_index:
        num_cols += 1

    table = shapes.add_table(num_rows, num_cols, left, top,
                             width, height).table

    if include_column_names:
        # write column headings
        for i, col in enumerate(data.columns):
            table.cell(0, i).text = str(col)

        row_shift = 1
    else:
        row_shift = 0

    # write row headings
    if include_index:
        for i, ind_val in enumerate(data.index):
            table.cell(i, 0).text = str(ind_val)

        col_shift = 1
    else:
        col_shift = 0

        # write body cells
    for i, col in enumerate(data.columns):
        for j in range(data.shape[0]):
            val = data.iloc[j, i]
            table.cell(j + row_shift, i + col_shift).text = str(val)

    return slide, table


def image_to_slide(presentation, img_path, slide_title="", left=None,
                   top=None, width=None, height=None):
    """ Add slide to specified presentation with provided image.

    Parameters
    ----------
    presentation : Presentation
        Presentation to add the slide to.

    img_path : str
        Path to the image file to open and create the image from.

    slide_title : str, optional
        Title of the slide, if any.

    left : Inches, float, int, optional
        Distance (in inches) from the left side of the new slide and the left
        side of the image.

    top : Inches, float, int, optional
        Distance (in inches) from the top of the new slide and the top of the
        image.

    width : Inches, float, int, optional
        Width, in inches, of the image to be inserted in the slide.

    height : Inches, float, int, optional
        Height, in inches, of the image to be inserted in the slide.
    """
    from pptx.presentation import Presentation

    if not isinstance(presentation, Presentation):
        msg = "The presentation should be a python-pptx Presentation."
        logger.exception(msg)
        raise ValueError(msg)

    title_only_slide_layout = presentation.slide_layouts[5]
    slide = presentation.slides.add_slide(title_only_slide_layout)

    shapes = slide.shapes
    if slide_title:
        title_shape = shapes.title
        title_shape.text = slide_title

    if left is None:
        left = Inches(1.8)
    elif isinstance(left, (float, int)):
        left = Inches(left)

    if top is None:
        top = Inches(1.8)
    elif isinstance(top, (float, int)):
        top = Inches(top)

    if height is None:
        height = Inches(5)
    elif isinstance(height, (float, int)):
        height = Inches(height)

    if width is None:
        width = Inches(6.5)
    elif isinstance(width, (float, int)):
        width = Inches(width)

    img = slide.shapes.add_picture(img_path, left=left, top=top, width=width,
                                   height=height)
    return img
