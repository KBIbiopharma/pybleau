from unittest import TestCase
from os.path import join, dirname
import pandas as pd
import pptx

from pybleau.reporting.pptx_utils import image_to_slide, Presentation, \
    title_slide, df_to_slide


HERE = dirname(__file__)


class TestGeneratePptx(TestCase):
    def setUp(self):
        self.presentation = Presentation()
        self.df = pd.DataFrame({"a": [1, 2], "b": [3, 4]})
        self.img_path = join(HERE, "test_img.png")

    def test_generate_pptx_with_title_slide(self):
        """ Test adding a title slide to existing presentation.
        """
        title_slide(presentation=self.presentation, title_text="Test title",
                    sub_title_text="Test subtitle")
        self.assertEqual(len(self.presentation.slides), 1)
        slide = self.presentation.slides[0]
        shapes = slide.shapes
        self.assertEqual(shapes.title.text, "Test title")
        self.assertEqual(len(shapes.placeholders), 2)

    def test_generate_pptx_with_title_slide2(self):
        """ Test making a new presentation when making a title slide.
        """
        pres = title_slide(title_text="Test title",
                           sub_title_text="Test subtitle")
        self.assertIsInstance(pres, pptx.presentation.Presentation)
        self.assertEqual(len(pres.slides), 1)
        slide = pres.slides[0]
        shapes = slide.shapes
        self.assertEqual(shapes.title.text, "Test title")
        self.assertEqual(len(shapes.placeholders), 2)

    def test_generate_pptx_with_data_table(self):
        """ Test adding a title slide to existing presentation.
        """
        for incl_col_header in [True, False]:
            for incl_row_header in [True, False]:
                self.presentation = Presentation()

                slide, table = df_to_slide(
                    self.presentation, self.df, slide_title="BLAH",
                    include_column_names=incl_col_header,
                    include_index=incl_row_header
                )
                self.assertEqual(len(self.presentation.slides), 1)
                slide2 = self.presentation.slides[0]
                self.assertIs(slide, slide2)
                shapes = slide.shapes
                self.assertEqual(shapes.title.text, "BLAH")
                self.assertEqual(len(table.columns), 2+int(incl_row_header))
                self.assertEqual(len(table.rows), 2+int(incl_col_header))

    def test_generate_pptx_with_images(self):
        img = image_to_slide(self.presentation, img_path=self.img_path,
                             slide_title="FOO")
        self.assertEqual(len(self.presentation.slides), 1)
        slide = self.presentation.slides[0]
        shapes = slide.shapes
        self.assertEqual(shapes.title.text, "FOO")
        self.assertIsInstance(img, pptx.shapes.picture.Picture)

    def test_add_all(self):
        title_slide(presentation=self.presentation, title_text="Test title",
                    sub_title_text="Test subtitle")
        df_to_slide(self.presentation, self.df, slide_title="BLAH")
        image_to_slide(self.presentation, img_path=self.img_path,
                       slide_title="FOO")
        self.assertEqual(len(self.presentation.slides), 3)
