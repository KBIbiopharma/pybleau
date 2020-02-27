""" Test exporting data and plots in a DFPlotManager
"""
from unittest import TestCase, skipIf
from pandas import DataFrame, read_csv, read_excel, read_hdf
import os
from shutil import rmtree
from os.path import dirname, isdir, isfile, join, splitext
import json
from zipfile import ZipFile
import pandas as pd

from pandas.testing import assert_frame_equal
from traits.testing.unittest_tools import UnittestTools

try:
    import kiwisolver  # noqa
    KIWI_AVAILABLE = True
except ImportError:
    KIWI_AVAILABLE = False

BACKEND_AVAILABLE = os.environ.get("ETS_TOOLKIT", "qt4") != "null"

if KIWI_AVAILABLE and BACKEND_AVAILABLE:
    from app_common.apptools.testing_utils import temp_bringup_ui_for
    from app_common.std_lib.filepath_utils import string2filename

    from pybleau.app.io.dataframe_plot_manager_exporter import \
        DataFramePlotManagerExporter, DEFAULT_DATASET_NAME, EXPORT_IN_FILE, \
        EXPORT_INLINE, EXPORT_NO, EXPORT_YES, IMG_FORMAT, PPT_FORMAT, \
        VEGA_FORMAT, EXTERNAL_DATA_FNAME
    from pybleau.app.model.dataframe_plot_manager import DataFramePlotManager
    from pybleau.reporting.string_definitions import IDX_NAME_KEY, \
        CONTENT_KEY, DATASETS_KEY
    from pybleau.app.model.plot_descriptor import PlotDescriptor
    from pybleau.app.plotting.api import BarPlotConfigurator, \
        HistogramPlotConfigurator, LinePlotConfigurator, \
        MultiHistogramPlotConfigurator, ScatterPlotConfigurator
    from pybleau.vega_translators.vega_chaco import df_to_vega, \
        TARGET_VEGA_SCHEMA

HERE = dirname(__file__)

TEST_DF = DataFrame({"a": [1, 2, 3, 4, 1, 2, 3, 4, 1, 2, 3, 4, 1, 2, 3, 4],
                     "b": [1, 1, 1, 1, 2, 2, 2, 2, 3, 3, 3, 3, 4, 4, 4, 4],
                     "c": [1, 2, 3, 4, 2, 3, 1, 1, 4, 4, 5, 6, 4, 4, 5, 6],
                     "d": list("ababcabcdabcdeab")},
                    dtype="float")

msg = "No UI backend to paint into or missing kiwisolver package"

if KIWI_AVAILABLE and BACKEND_AVAILABLE:
    class NonInteractiveExporter(DataFramePlotManagerExporter):
        interactive = False


@skipIf(not BACKEND_AVAILABLE or not KIWI_AVAILABLE, msg)
class TestPlotManagerExporterView(TestCase, UnittestTools):
    def setUp(self):
        self.model = DataFramePlotManager(data_source=TEST_DF)

        self.exporter = NonInteractiveExporter(
            df_plotter=self.model, target_dir=".",
        )

    def test_bring_up(self):
        exporter = self.exporter
        with temp_bringup_ui_for(exporter):
            for export_format in {IMG_FORMAT, VEGA_FORMAT, PPT_FORMAT}:
                exporter.export_format = export_format

    def test_update_export_data_options(self):
        exporter = self.exporter
        exporter.export_format = IMG_FORMAT
        with self.assertTraitChanges(exporter, "_export_data_options"):
            exporter.export_format = VEGA_FORMAT

        with self.assertTraitChanges(exporter, "_export_data_options"):
            exporter.export_format = PPT_FORMAT


class BasePlotManagerExportTester(object):
    def setUp(self):
        self.config = HistogramPlotConfigurator(data_source=TEST_DF)
        self.config.x_col_name = "a"
        self.desc = PlotDescriptor(x_col_name="a", plot_config=self.config,
                                   plot_title="Plot 0")

        self.config2 = BarPlotConfigurator(data_source=TEST_DF)
        self.config2.x_col_name = "a"
        self.config2.y_col_name = "b"
        self.desc2 = PlotDescriptor(x_col_name="a", y_col_name="b",
                                    plot_config=self.config2,
                                    plot_title="Plot 1")

        self.config3 = ScatterPlotConfigurator(data_source=TEST_DF)
        self.config3.x_col_name = "a"
        self.config3.y_col_name = "b"
        self.desc3 = PlotDescriptor(x_col_name="a", y_col_name="b",
                                    plot_config=self.config3,
                                    plot_title="Plot 2")

        self.config4 = ScatterPlotConfigurator(data_source=TEST_DF)
        self.config4.x_col_name = "a"
        self.config4.y_col_name = "b"
        self.config4.z_col_name = "d"
        self.desc4 = PlotDescriptor(x_col_name="a", y_col_name="b",
                                    z_col_name="d", plot_config=self.config4,
                                    plot_title="Plot 3")

        self.model = DataFramePlotManager(contained_plots=[self.desc],
                                          data_source=TEST_DF)
        self.target_dir = join(HERE, "test_dir")

        if isdir(self.target_dir):
            rmtree(self.target_dir)

    def tearDown(self):
        if isdir(self.target_dir):
            rmtree(self.target_dir)

    # Tests -------------------------------------------------------------------

    def test_overwrite_flag(self):
        model = DataFramePlotManager(contained_plots=[self.desc],
                                     data_source=TEST_DF)
        exporter = NonInteractiveExporter(
            df_plotter=model, target_dir=self.target_dir,
            target_file=self.target_file
        )
        export_func = getattr(exporter, self.converter)
        export_func()

        self.assertIn(self.target_filename, os.listdir(self.target_dir))

        # Can't do it a second time because the file is already there...
        with self.assertRaises(IOError):
            export_func()

        # ... unless the overwrite flag is True:
        exporter = NonInteractiveExporter(
            df_plotter=model, target_dir=self.target_dir,
            overwrite_file_if_exists=True
        )
        export_func = getattr(exporter, self.converter)
        export_func()

    # Assertion methods -------------------------------------------------------

    def assert_data_in_file(self, data_file, num_dataset=1):
        """ Check that the target data file has the right content.
        """
        ext = splitext(data_file)[1]
        kw = {}
        if ext == ".xlsx":
            reader = read_excel
            kw["sheet_name"] = DEFAULT_DATASET_NAME
            kw["index_col"] = 0
        elif ext == ".csv":
            reader = read_csv
            kw["index_col"] = 0
        elif ext == ".h5":
            reader = read_hdf
            kw["key"] = DEFAULT_DATASET_NAME

        df_back = reader(data_file, **kw)
        # FIXME: is there any way to avoid Excel messing up the dtypes?
        if ext == ".xlsx":
            assert_frame_equal(df_back, TEST_DF, check_dtype=False)
        else:
            assert_frame_equal(df_back, TEST_DF)

        if num_dataset > 1:
            if ext == ".h5":
                with pd.HDFStore(data_file) as s:
                    self.assertEqual(len(s.keys()), num_dataset)
                    for key in s.keys():
                        self.assertIsInstance(s[key], pd.DataFrame)
            elif ext == ".xlsx":
                with pd.ExcelFile(data_file) as f:
                    self.assertEqual(len(f.sheet_names), num_dataset)
                    for key in f.sheet_names:
                        self.assertIsInstance(reader(f, sheet_name=key),
                                              pd.DataFrame)
            elif ext == ".csv":
                content = os.listdir(dirname(data_file))
                self.assertEqual(len(content), num_dataset)
                for fname in content:
                    df = read_csv(join(dirname(data_file), fname))
                    self.assertIsInstance(df, pd.DataFrame)


@skipIf(not BACKEND_AVAILABLE or not KIWI_AVAILABLE, msg)
class TestPlotManagerExportImages(BasePlotManagerExportTester, TestCase):
    def setUp(self):
        super(TestPlotManagerExportImages, self).setUp()
        self.converter = "to_folder"
        self.target_filename = "0_Plot_0.PNG"
        self.target_file = join(self.target_dir, self.target_filename)

    def test_export_1_plot(self):
        exporter = NonInteractiveExporter(df_plotter=self.model,
                                          target_dir=self.target_dir)
        export_func = getattr(exporter, self.converter)
        export_func()
        content = os.listdir(self.target_dir)
        self.assertEqual(len(content), 1)
        fname = content[0]
        self.assertEqual(fname, self.target_filename)

    def test_export_1_plot_special_title(self):
        self.desc.plot_title = "blah!"
        model = DataFramePlotManager(contained_plots=[self.desc],
                                     data_source=TEST_DF)
        exporter = NonInteractiveExporter(df_plotter=model,
                                          target_dir=self.target_dir)
        export_func = getattr(exporter, self.converter)
        export_func()
        content = os.listdir(self.target_dir)
        self.assertEqual(len(content), 1)
        fname = content[0]
        self.assertEqual(fname, "0_blah.PNG")

    def test_export_1_plot_in_jpeg(self):
        exporter = NonInteractiveExporter(df_plotter=self.model,
                                          target_dir=self.target_dir,
                                          image_format="JPG")
        export_func = getattr(exporter, self.converter)
        export_func()
        content = os.listdir(self.target_dir)
        self.assertEqual(len(content), 1)
        fname = content[0]
        self.assertEqual(fname, "0_Plot_0.JPG")

    def test_export_3_plots_no_data(self):
        model = DataFramePlotManager(contained_plots=[self.desc, self.desc2,
                                                      self.desc3],
                                     data_source=TEST_DF)
        exporter = NonInteractiveExporter(df_plotter=model,
                                          target_dir=self.target_dir)
        export_func = getattr(exporter, self.converter)
        export_func()
        content = os.listdir(self.target_dir)
        self.assertEqual(len(content), 3)
        self.assertEqual(set(content), {"0_Plot_0.PNG", "1_Plot_1.PNG",
                                        "2_Plot_2.PNG"})

    def test_export_1_plot_export_source_data(self):
        for fmt in [".h5", ".csv", ".xlsx"]:
            exporter = NonInteractiveExporter(df_plotter=self.model,
                                              target_dir=self.target_dir,
                                              export_data=EXPORT_YES,
                                              data_format=fmt)
            export_func = getattr(exporter, self.converter)
            export_func()
            content = os.listdir(self.target_dir)
            self.assertEqual(len(content), 2)
            self.assertEqual(set(content),
                             {self.target_filename, EXTERNAL_DATA_FNAME+fmt})

            # Check data file content:
            data_file = join(self.target_dir, EXTERNAL_DATA_FNAME+fmt)
            self.assert_data_in_file(data_file)

            self.tearDown()

    def test_export_3_plots_export_source_data(self):
        model = DataFramePlotManager(contained_plots=[self.desc, self.desc2,
                                                      self.desc3],
                                     data_source=TEST_DF)

        for fmt in [".h5", ".csv", ".xlsx"]:
            exporter = NonInteractiveExporter(df_plotter=model,
                                              target_dir=self.target_dir,
                                              export_data=EXPORT_YES,
                                              data_format=fmt)
            export_func = getattr(exporter, self.converter)
            export_func()
            content = os.listdir(self.target_dir)
            self.assertEqual(len(content), 4)
            expected = {"0_Plot_0.PNG", "1_Plot_1.PNG", "2_Plot_2.PNG",
                        EXTERNAL_DATA_FNAME + fmt}
            self.assertEqual(set(content), expected)

            # Check data file content: csv isn't zipped because there is only 1
            # file:
            data_file = join(self.target_dir, EXTERNAL_DATA_FNAME + fmt)
            self.assert_data_in_file(data_file)

            self.tearDown()

    def test_export_multi_renderer_plot_export_source_and_plot_data(self):
        model = DataFramePlotManager(contained_plots=[self.desc4],
                                     data_source=TEST_DF)
        for fmt in [".h5", ".csv", ".xlsx"]:
            exporter = NonInteractiveExporter(
                df_plotter=model, target_dir=self.target_dir,
                export_data=EXPORT_YES, data_format=fmt,
                export_each_plot_data=True
            )
            export_func = getattr(exporter, self.converter)
            export_func()
            content = os.listdir(self.target_dir)
            self.assertEqual(len(content), 2)
            if fmt in [".h5", ".xlsx"]:
                expected = {"0_Plot_3.PNG", EXTERNAL_DATA_FNAME + fmt}
            else:
                expected = {"0_Plot_3.PNG", EXTERNAL_DATA_FNAME + ".zip"}

            self.assertEqual(set(content), expected)

            # Check data file content:
            if fmt == ".csv":
                zip_path = join(self.target_dir, EXTERNAL_DATA_FNAME + ".zip")
                with ZipFile(zip_path) as f:
                    f.extractall(join(self.target_dir, "data"))

                source_fname = string2filename(DEFAULT_DATASET_NAME) + ".csv"
                data_file = join(self.target_dir, "data", source_fname)
            else:
                data_file = join(self.target_dir, EXTERNAL_DATA_FNAME + fmt)

            num_d = len(set(TEST_DF["d"]))
            self.assert_data_in_file(data_file, num_dataset=1+num_d)

            # clean up at every loop:
            self.tearDown()


@skipIf(not BACKEND_AVAILABLE or not KIWI_AVAILABLE, msg)
class TestPlotManagerExportPptx(BasePlotManagerExportTester, TestCase):

    def setUp(self):
        super(TestPlotManagerExportPptx, self).setUp()
        self.converter = "to_pptx"
        self.target_filename = "test.pptx"
        self.target_file = join(self.target_dir, self.target_filename)

    def test_export_1_plot(self):
        model = DataFramePlotManager(contained_plots=[self.desc],
                                     data_source=TEST_DF)
        exporter = NonInteractiveExporter(df_plotter=model,
                                          target_file=self.target_file)
        export_func = getattr(exporter, self.converter)
        export_func()
        content = os.listdir(self.target_dir)
        self.assertEqual(len(content), 1)
        fname = content[0]
        self.assertEqual(fname, self.target_filename)
        self.assert_valid_pptx(exporter)

    def test_export_1_plot_export_source_data(self):
        model = DataFramePlotManager(contained_plots=[self.desc],
                                     data_source=TEST_DF)
        for fmt in [".h5", ".csv", ".xlsx"]:
            exporter = NonInteractiveExporter(
                df_plotter=model, target_file=self.target_file,
                export_data=EXPORT_YES, data_format=fmt
            )
            export_func = getattr(exporter, self.converter)
            export_func()
            content = os.listdir(self.target_dir)
            self.assertEqual(len(content), 2)
            self.assertEqual(set(content),
                             {EXTERNAL_DATA_FNAME + fmt, self.target_filename})
            self.assert_valid_pptx(exporter)
            fname = join(self.target_dir, EXTERNAL_DATA_FNAME + fmt)
            self.assert_data_in_file(fname)

            # clean up at every loop:
            self.tearDown()

    def test_export_1_plot_with_custom_title(self):
        model = DataFramePlotManager(contained_plots=[self.desc],
                                     data_source=TEST_DF)
        exporter = NonInteractiveExporter(df_plotter=model,
                                          target_file=self.target_file,
                                          presentation_title="blah")
        export_func = getattr(exporter, self.converter)
        export_func()
        content = os.listdir(self.target_dir)
        self.assertEqual(len(content), 1)
        fname = content[0]
        self.assertEqual(fname, self.target_filename)
        self.assert_valid_pptx(exporter, title_text="blah")

    def test_export_3_plots_no_data(self):

        model = DataFramePlotManager(contained_plots=[self.desc, self.desc2,
                                                      self.desc3],
                                     data_source=TEST_DF)
        exporter = NonInteractiveExporter(df_plotter=model,
                                          target_file=self.target_file)
        export_func = getattr(exporter, self.converter)
        export_func()
        content = os.listdir(self.target_dir)
        self.assertEqual(len(content), 1)
        self.assertEqual(content[0], self.target_filename)
        self.assert_valid_pptx(exporter, num_images=3)

    def test_export_3_plots_export_source_data(self):
        model = DataFramePlotManager(contained_plots=[self.desc, self.desc2,
                                                      self.desc3],
                                     data_source=TEST_DF)
        for fmt in [".h5", ".csv", ".xlsx"]:
            exporter = NonInteractiveExporter(
                df_plotter=model, target_file=self.target_file,
                export_data=EXPORT_YES, data_format=fmt,
            )
            export_func = getattr(exporter, self.converter)
            export_func()
            content = os.listdir(self.target_dir)
            self.assertEqual(len(content), 2)
            expected = {self.target_filename, EXTERNAL_DATA_FNAME + fmt}
            self.assertEqual(set(content), expected)

            # Check data file content:
            data_file = join(self.target_dir, EXTERNAL_DATA_FNAME + fmt)
            self.assert_data_in_file(data_file)

            # clean up at every loop:
            self.tearDown()

    def test_export_3_plots_export_source_and_plot_data(self):
        model = DataFramePlotManager(contained_plots=[self.desc, self.desc2,
                                                      self.desc3],
                                     data_source=TEST_DF)
        for fmt in [".h5", ".csv", ".xlsx"]:
            exporter = NonInteractiveExporter(
                df_plotter=model, target_file=self.target_file,
                export_data=EXPORT_YES, data_format=fmt,
                export_each_plot_data=True
            )
            export_func = getattr(exporter, self.converter)
            export_func()
            content = os.listdir(self.target_dir)
            self.assertEqual(len(content), 2)
            if fmt in [".h5", ".xlsx"]:
                expected = {self.target_filename, EXTERNAL_DATA_FNAME + fmt}
            else:
                expected = {self.target_filename, EXTERNAL_DATA_FNAME + ".zip"}

            self.assertEqual(set(content), expected)

            # Check data file content:
            if fmt == ".csv":
                zip_path = join(self.target_dir, EXTERNAL_DATA_FNAME + ".zip")
                with ZipFile(zip_path) as f:
                    f.extractall(join(self.target_dir, "data"))

                source_fname = string2filename(DEFAULT_DATASET_NAME) + ".csv"
                data_file = join(self.target_dir, "data", source_fname)
            else:
                data_file = join(self.target_dir, EXTERNAL_DATA_FNAME + fmt)

            self.assert_data_in_file(data_file)

            # clean up at every loop:
            self.tearDown()

    def test_export_multi_renderer_plot_export_source_and_plot_data(self):
        model = DataFramePlotManager(contained_plots=[self.desc4],
                                     data_source=TEST_DF)
        for fmt in [".h5", ".csv", ".xlsx"]:
            exporter = NonInteractiveExporter(
                df_plotter=model, target_file=self.target_file,
                export_data=EXPORT_YES, data_format=fmt,
                export_each_plot_data=True
            )
            export_func = getattr(exporter, self.converter)
            export_func()
            content = os.listdir(self.target_dir)
            self.assertEqual(len(content), 2)
            if fmt in [".h5", ".xlsx"]:
                expected = {self.target_filename, EXTERNAL_DATA_FNAME + fmt}
            else:
                expected = {self.target_filename, EXTERNAL_DATA_FNAME + ".zip"}

            self.assertEqual(set(content), expected)

            self.assert_valid_pptx(exporter)

            # Check data file content:
            if fmt == ".csv":
                zip_path = join(self.target_dir, EXTERNAL_DATA_FNAME + ".zip")
                with ZipFile(zip_path) as f:
                    f.extractall(join(self.target_dir, "data"))

                source_fname = string2filename(DEFAULT_DATASET_NAME) + ".csv"
                data_file = join(self.target_dir, "data", source_fname)
            else:
                data_file = join(self.target_dir, EXTERNAL_DATA_FNAME + fmt)

            num_d = len(set(TEST_DF["d"]))
            self.assert_data_in_file(data_file, num_dataset=1 + num_d)

            # clean up at every loop:
            self.tearDown()

    # Utility method ----------------------------------------------------------

    def assert_valid_pptx(self, exporter, num_images=1, title_text=None):
        from pptx import Presentation
        from pptx.shapes.picture import Picture

        if title_text is None:
            title_text = 'New Presentation'

        prs = Presentation(pptx=exporter.target_file)
        # Title slide and 1 image slide
        self.assertEqual(len(list(prs.slides)), num_images+1)
        title_slide = list(prs.slides)[0]
        self.assertEqual(title_slide.shapes[0].text, title_text)
        for i, slide in enumerate(list(prs.slides)[1:]):
            slide_title, slide_img = list(slide.shapes)
            self.assertEqual(slide_title.text, "plot_{}".format(i))
            self.assertIsInstance(slide_img, Picture)


@skipIf(not BACKEND_AVAILABLE or not KIWI_AVAILABLE, msg)
class TestPlotManagerExportVega(BasePlotManagerExportTester, TestCase,
                                UnittestTools):
    def setUp(self):
        super(TestPlotManagerExportVega, self).setUp()

        if isdir(self.target_dir):
            rmtree(self.target_dir)

        self.converter = "to_vega"
        self.model = DataFramePlotManager(data_source=TEST_DF)
        self.exporter = NonInteractiveExporter(
            df_plotter=self.model, target_dir=self.target_dir,
            export_format=VEGA_FORMAT, export_data=EXPORT_NO
        )
        self.target_filename = "temp.json"
        self.target_file = join(self.target_dir, self.target_filename)

    def test_export_no_plot_no_data_no_file(self):
        content = self.exporter.to_vega()
        self.assertEqual(content[CONTENT_KEY], [])
        self.assertEqual(content[DATASETS_KEY], {})

    def test_export_no_plot_w_data_no_file(self):
        self.exporter.export_data = EXPORT_IN_FILE
        content = self.exporter.to_vega()
        self.assertEqual(content[CONTENT_KEY], [])
        self.assert_rebuild_df(content)

    def test_export_line_plot_w_data(self):
        config = LinePlotConfigurator(data_source=TEST_DF,
                                      plot_title="Plot")
        config.x_col_name = "a"
        config.y_col_name = "b"
        self.model._add_new_plot(config)

        self.exporter.export_data = EXPORT_IN_FILE
        content = self.exporter.to_vega()
        dataset = content[DATASETS_KEY][DEFAULT_DATASET_NAME]
        self.assertIsInstance(dataset["data"], list)
        self.assert_rebuild_df(content, cascading_data=False)
        self.assertEqual(len(content[CONTENT_KEY]), 1)
        expected = {'$schema': TARGET_VEGA_SCHEMA,
                    'data': {"name": DEFAULT_DATASET_NAME},
                    'encoding': {'x': {'field': 'a', 'type': 'quantitative'},
                                 'y': {'field': 'b', 'type': 'quantitative'}},
                    'mark': 'line'}
        self.assert_vega_export_equal(content[CONTENT_KEY][0], expected)

    def test_export_1_plot_to_file(self):
        config = LinePlotConfigurator(data_source=TEST_DF,
                                      plot_title="Plot")
        config.x_col_name = "a"
        config.y_col_name = "b"
        self.model._add_new_plot(config)

        self.exporter.target_file = self.target_file
        content = self.exporter.to_vega()
        self.assertTrue(isfile(self.target_file))
        file_content = json.load(open(self.target_file))
        self.assertEqual(content, file_content)
        expected = {'$schema': TARGET_VEGA_SCHEMA,
                    'data': {},
                    'encoding': {'x': {'field': 'a', 'type': 'quantitative'},
                                 'y': {'field': 'b', 'type': 'quantitative'}},
                    'mark': 'line'}
        self.assert_vega_export_equal(content[CONTENT_KEY][0], expected)
        self.assertEqual(content[DATASETS_KEY], {})

    def test_export_2line_plots_w_data(self):
        config = LinePlotConfigurator(data_source=TEST_DF,
                                      plot_title="Plot")
        config.x_col_name = "a"
        config.y_col_name = "b"
        self.model._add_new_plot(config)

        config = LinePlotConfigurator(data_source=TEST_DF,
                                      plot_title="Plot2")
        config.x_col_name = "b"
        config.y_col_name = "a"
        self.model._add_new_plot(config)

        self.exporter.export_data = EXPORT_IN_FILE
        content = self.exporter.to_vega()
        self.assert_rebuild_df(content)
        self.assertEqual(len(content[CONTENT_KEY]), 2)
        expected = {'$schema': TARGET_VEGA_SCHEMA,
                    'data': {"name": DEFAULT_DATASET_NAME},
                    'encoding': {'x': {'field': 'a', 'type': 'quantitative'},
                                 'y': {'field': 'b', 'type': 'quantitative'}},
                    'mark': 'line'}
        self.assert_vega_export_equal(content[CONTENT_KEY][0], expected)

        expected = {'$schema': TARGET_VEGA_SCHEMA,
                    'data': {"name": DEFAULT_DATASET_NAME},
                    'encoding': {'x': {'field': 'b', 'type': 'quantitative'},
                                 'y': {'field': 'a', 'type': 'quantitative'}},
                    'mark': 'line'}
        self.assert_vega_export_equal(content[CONTENT_KEY][1], expected)

    def test_export_2line_plots_w_cascading_data(self):
        config = LinePlotConfigurator(data_source=TEST_DF,
                                      plot_title="Plot")
        config.x_col_name = "a"
        config.y_col_name = "b"
        self.model._add_new_plot(config)

        config = LinePlotConfigurator(data_source=TEST_DF,
                                      plot_title="Plot2")
        config.x_col_name = "b"
        config.y_col_name = "a"
        self.model._add_new_plot(config)

        # In the cascade style, the data is rewritten in the description of
        # every plot:
        self.exporter.export_data = EXPORT_INLINE
        content = self.exporter.to_vega()
        self.assertEqual(content[DATASETS_KEY], {})
        self.assert_rebuild_df(content, cascading_data=True)
        self.assertEqual(len(content[CONTENT_KEY]), 2)
        expected = {'$schema': TARGET_VEGA_SCHEMA,
                    'data': {"values": df_to_vega(TEST_DF),
                             IDX_NAME_KEY: None},
                    'encoding': {'x': {'field': 'a', 'type': 'quantitative'},
                                 'y': {'field': 'b', 'type': 'quantitative'}},
                    'mark': 'line'}
        self.assert_vega_export_equal(content[CONTENT_KEY][0], expected)

        expected = {'$schema': TARGET_VEGA_SCHEMA,
                    'data': {"values": df_to_vega(TEST_DF),
                             IDX_NAME_KEY: None},
                    'encoding': {'x': {'field': 'b', 'type': 'quantitative'},
                                 'y': {'field': 'a', 'type': 'quantitative'}},
                    'mark': 'line'}
        self.assert_vega_export_equal(content[CONTENT_KEY][1], expected)

    def test_export_2scatter_plots_w_data(self):
        config = ScatterPlotConfigurator(data_source=TEST_DF,
                                         plot_title="Plot")
        config.x_col_name = "a"
        config.y_col_name = "b"
        self.model._add_new_plot(config)

        config = ScatterPlotConfigurator(data_source=TEST_DF,
                                         plot_title="Plot2")
        config.x_col_name = "b"
        config.y_col_name = "a"
        self.model._add_new_plot(config)

        self.exporter.export_data = EXPORT_IN_FILE
        content = self.exporter.to_vega()
        self.assert_rebuild_df(content)
        self.assertEqual(len(content[CONTENT_KEY]), 2)
        expected = {'$schema': TARGET_VEGA_SCHEMA,
                    'data': {"name": DEFAULT_DATASET_NAME},
                    'encoding': {'x': {'field': 'a', 'type': 'quantitative'},
                                 'y': {'field': 'b', 'type': 'quantitative'}},
                    'mark': 'point'}
        self.assert_vega_export_equal(content[CONTENT_KEY][0], expected)

        expected = {'$schema': TARGET_VEGA_SCHEMA,
                    'data': {"name": DEFAULT_DATASET_NAME},
                    'encoding': {'x': {'field': 'b', 'type': 'quantitative'},
                                 'y': {'field': 'a', 'type': 'quantitative'}},
                    'mark': 'point'}
        self.assert_vega_export_equal(content[CONTENT_KEY][1], expected)

    def test_export_2hist_plots_w_data(self):
        config = HistogramPlotConfigurator(data_source=TEST_DF,
                                           plot_title="Plot")
        config.x_col_name = "a"
        self.model._add_new_plot(config)

        config = HistogramPlotConfigurator(data_source=TEST_DF,
                                           plot_title="Plot2")
        config.x_col_name = "b"
        self.model._add_new_plot(config)

        self.exporter.export_data = EXPORT_IN_FILE
        content = self.exporter.to_vega()
        self.assert_rebuild_df(content)
        self.assertEqual(len(content[CONTENT_KEY]), 2)
        expected = {'$schema': TARGET_VEGA_SCHEMA,
                    'data': {"name": DEFAULT_DATASET_NAME},
                    'encoding': {
                        "x": {
                          "bin": True,
                          "field": "a",
                          "type": "quantitative"
                        },
                        "y": {
                          "aggregate": "count",
                          "type": "quantitative"
                        }
                    },
                    'mark': 'bar'}
        self.assert_vega_export_equal(content[CONTENT_KEY][0], expected)

        expected = {'$schema': TARGET_VEGA_SCHEMA,
                    'data': {"name": DEFAULT_DATASET_NAME},
                    'encoding': {
                        "x": {
                          "bin": True,
                          "field": "b",
                          "type": "quantitative"
                        },
                        "y": {
                          "aggregate": "count",
                          "type": "quantitative"
                        }
                    },
                    'mark': 'bar'}
        self.assert_vega_export_equal(content[CONTENT_KEY][1], expected)

    def test_export_multi_hist_plots_w_data(self):
        config = MultiHistogramPlotConfigurator(data_source=TEST_DF,
                                                plot_title="Plot {i}")
        config.x_col_names = ["a", "b"]
        self.model._add_new_plots(config)

        self.exporter.export_data = EXPORT_IN_FILE
        content = self.exporter.to_vega()
        self.assert_rebuild_df(content)
        self.assertEqual(len(content[CONTENT_KEY]), 2)
        expected = {'$schema': TARGET_VEGA_SCHEMA,
                    'data': {"name": DEFAULT_DATASET_NAME},
                    'encoding': {
                        "x": {
                          "bin": True,
                          "field": "a",
                          "type": "quantitative"
                        },
                        "y": {
                          "aggregate": "count",
                          "type": "quantitative"
                        }
                    },
                    'mark': 'bar'}
        self.assert_vega_export_equal(content[CONTENT_KEY][0], expected)

        expected = {'$schema': TARGET_VEGA_SCHEMA,
                    'data': {"name": DEFAULT_DATASET_NAME},
                    'encoding': {
                        "x": {
                          "bin": True,
                          "field": "b",
                          "type": "quantitative"
                        },
                        "y": {
                          "aggregate": "count",
                          "type": "quantitative"
                        }
                    },
                    'mark': 'bar'}
        self.assert_vega_export_equal(content[CONTENT_KEY][1], expected)

    # Assertion methods -------------------------------------------------------

    def assert_vega_export_equal(self, desc1, desc2):
        self.assertEqual(set(desc1.keys()), set(desc2.keys()))
        for key in desc1:
            self.assertEqual(desc1[key], desc2[key],
                             msg="{} unequal".format(key))

    def assert_rebuild_df(self, content, cascading_data=False):
        if cascading_data:
            for desc in content[CONTENT_KEY]:
                dataset = desc["data"]
                idx_col = dataset[IDX_NAME_KEY]
                if idx_col is None:
                    rebuilt_df = DataFrame(dataset["values"]).set_index(
                        "index")
                    rebuilt_df.index.name = None
                else:
                    rebuilt_df = DataFrame(dataset["values"]).set_index(
                        idx_col)

                assert_frame_equal(rebuilt_df, TEST_DF, check_index_type=False,
                                   check_names=False)
        else:
            datasets = content[DATASETS_KEY]
            for name, dataset in datasets.items():
                self.assertIsInstance(dataset["data"], list)
                idx_col = dataset[IDX_NAME_KEY]
                if idx_col is None:
                    rebuilt_df = DataFrame(dataset["data"]).set_index("index")
                    rebuilt_df.index.name = None
                else:
                    rebuilt_df = DataFrame(dataset["data"]).set_index(idx_col)

                assert_frame_equal(rebuilt_df, TEST_DF, check_index_type=False,
                                   check_names=False)
